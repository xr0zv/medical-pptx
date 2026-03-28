#!/usr/bin/env python3
"""
Figure Insertion Script for medical-pptx presentations.

Scans a PPTX file for placeholder shapes matching patterns like
  [Insert Fig. X — description]  or  [Figure: Fig. X — description]
and replaces them with the corresponding extracted figure images,
using manifest.json metadata for matching and aspect-ratio-correct sizing.
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# Placeholder detection
# ---------------------------------------------------------------------------

# Matches:  [Insert Fig. 35.1 — Lumbosacral plexus anatomy]
#           [Insert Fig. 35.1 - Lumbosacral plexus anatomy]
#           [Figure: Fig. 35.1 — Lumbosacral plexus anatomy]
#           [Figure: Table 1 - Something]
#           [Insert CASE 1 — Clinical scenario]
#           [Insert Image 1 — Radiograph]
#           [Insert Plate 2 — Histology]
#           [Insert Panel A — Subpanel description]
PLACEHOLDER_RE = re.compile(
    r"\[\s*(?:Insert|Figure:)\s+"    # opening bracket + keyword
    r"((?:Fig\.?|Figure|Table|Box|CASE|Image|Plate)\s*[\d]+(?:[\.\-]\d+)*"  # label: most types require digits
    r"|Panel\s*[A-Za-z\d]+)"        # Panel allows letter suffixes (Panel A, Panel B)
    r"\s*[\u2014\-]\s*"             # em-dash or hyphen separator
    r"(.+?)"                        # description (not used for matching)
    r"\s*\]",                        # closing bracket
    re.IGNORECASE,
)


def normalize_label(label: str) -> str:
    """Lowercase, collapse whitespace, normalize dot-digit spacing, strip trailing dots."""
    s = re.sub(r"\s+", " ", label.strip().lower())
    # Normalize "fig.7" → "fig. 7" but NOT "35.1" → "35. 1"
    # Only insert space when a letter precedes the dot (label prefix boundary)
    s = re.sub(r"([a-z])\.(\d)", r"\1. \2", s)
    return s.rstrip(".")


def extract_placeholder_label(text: str) -> Optional[str]:
    """Return the figure label from placeholder text, or None."""
    m = PLACEHOLDER_RE.search(text)
    if m:
        return m.group(1).strip()
    return None


# ---------------------------------------------------------------------------
# Aspect-ratio-correct fitting (mirrors fitImage from the skill)
# ---------------------------------------------------------------------------

def fit_image(
    img_aspect: float,
    box_left: int,
    box_top: int,
    box_width: int,
    box_height: int,
) -> Tuple[int, int, int, int]:
    """Compute (left, top, width, height) in EMUs to fit an image inside a box.

    The image is scaled to fit within box_width x box_height while preserving
    its aspect ratio, then centered within the bounding box.
    """
    disp_w = box_width
    disp_h = int(disp_w / img_aspect)

    if disp_h > box_height:
        disp_h = box_height
        disp_w = int(disp_h * img_aspect)

    # Center within the bounding box
    x = box_left + (box_width - disp_w) // 2
    y = box_top + (box_height - disp_h) // 2

    return x, y, disp_w, disp_h


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def shape_text(shape) -> str:
    """Safely extract text from a shape."""
    if shape.has_text_frame:
        return shape.text_frame.text
    return ""


def find_companion_rectangle(slide, text_shape) -> Optional[object]:
    """Try to find the dashed-border rectangle behind a placeholder text box.

    Heuristic: look for a rectangle shape on the same slide that fully
    contains the text shape and has no text frame (or empty text).  Among
    candidates, pick the one whose area is closest to the text shape area
    (likely the deliberately-created backdrop).
    """
    ts_left = text_shape.left
    ts_top = text_shape.top
    ts_right = ts_left + text_shape.width
    ts_bottom = ts_top + text_shape.height

    best = None
    best_area_diff = None

    for s in slide.shapes:
        if s.shape_id == text_shape.shape_id:
            continue
        # Must be a rectangle-like shape with no meaningful text
        if shape_text(s).strip():
            continue
        # Check containment: candidate must enclose or roughly overlap the text
        s_left = s.left
        s_top = s.top
        s_right = s_left + s.width
        s_bottom = s_top + s.height

        # The rectangle should contain the text box (with a small tolerance)
        tolerance = Emu(45720)  # ~0.05 inches
        if (s_left <= ts_left + tolerance and
                s_top <= ts_top + tolerance and
                s_right >= ts_right - tolerance and
                s_bottom >= ts_bottom - tolerance):
            area_diff = abs(s.width * s.height - text_shape.width * text_shape.height)
            if best is None or area_diff < best_area_diff:
                best = s
                best_area_diff = area_diff

    return best


# ---------------------------------------------------------------------------
# Main logic
# ---------------------------------------------------------------------------

def load_manifest(manifest_path: str) -> List[dict]:
    with open(manifest_path, "r", encoding="utf-8") as f:
        return json.load(f)


def build_label_index(manifest: List[dict]) -> Dict[str, dict]:
    """Map normalized labels to manifest entries."""
    index: Dict[str, dict] = {}
    for entry in manifest:
        key = normalize_label(entry.get("label", ""))
        if key:
            index[key] = entry
    return index


def process_presentation(
    pptx_path: str,
    manifest_path: str,
    output_path: str,
    debug: bool = False,
) -> bool:
    """Insert figures into placeholders. Returns True on success."""

    manifest = load_manifest(manifest_path)
    label_index = build_label_index(manifest)
    manifest_dir = Path(manifest_path).parent

    prs = Presentation(pptx_path)

    matched_labels: List[str] = []
    unmatched_placeholders: List[str] = []
    used_labels: set = set()

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Collect shapes first to avoid mutating while iterating
        shapes_snapshot = list(slide.shapes)

        for shape in shapes_snapshot:
            text = shape_text(shape)
            if not text:
                continue

            label = extract_placeholder_label(text)
            if label is None:
                continue

            norm = normalize_label(label)
            entry = label_index.get(norm)

            if entry is None:
                msg = f"  Slide {slide_idx}: placeholder '{label}' — no matching figure in manifest"
                print(msg, file=sys.stderr)
                unmatched_placeholders.append(label)
                continue

            # Resolve image path relative to manifest directory
            img_rel = entry["file"]
            img_path = str(manifest_dir / img_rel)
            if not os.path.isfile(img_path):
                msg = f"  Slide {slide_idx}: placeholder '{label}' — image file not found: {img_path}"
                print(msg, file=sys.stderr)
                unmatched_placeholders.append(label)
                continue

            # Determine bounding box.
            # The text box is sized to match the companion rectangle (same x/y/w/h),
            # so its own dimensions are always correct for image placement.
            # Companion detection is only used for cleanup (removing the dashed rect).
            box_left = shape.left
            box_top = shape.top
            box_width = shape.width
            box_height = shape.height
            companion = find_companion_rectangle(slide, shape)

            # Compute aspect-ratio-correct placement
            aspect = entry.get("aspect_ratio")
            if aspect is None:
                w = entry.get("width", 1)
                h = entry.get("height", 1)
                aspect = w / h if h else 1.0

            x, y, w, h = fit_image(aspect, box_left, box_top, box_width, box_height)

            # Add image
            slide.shapes.add_picture(img_path, x, y, w, h)

            # Reposition placeholder text box as a reference label below the image
            SLIDE_HEIGHT = Emu(5143500)  # 5.625 inches (16:9 slide height)
            label_height = Emu(274320)   # ~0.3 inches

            # Clip label to slide bottom edge if image sits low on the slide
            label_top = y + h
            if label_top + label_height > SLIDE_HEIGHT:
                label_top = SLIDE_HEIGHT - label_height

            shape.left = x
            shape.top = label_top
            shape.width = w
            shape.height = label_height

            # Clean up text: "[Insert Fig. 35.1 — desc]" → "Fig. 35.1 — desc"
            raw = shape.text_frame.text
            cleaned = re.sub(r"^\[\s*(?:Insert|Figure:)\s+", "", raw)
            cleaned = re.sub(r"\s*\]$", "", cleaned)

            # Set label text with explicit formatting (don't rely on inherited styles)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = cleaned
                    para.runs[0].font.size = Pt(10)
                    para.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
                    para.runs[0].font.italic = True
                    para.runs[0].font.name = "Calibri"
                else:
                    # No runs exist — add one explicitly
                    run = para.add_run()
                    run.text = cleaned
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
                    run.font.italic = True
                    run.font.name = "Calibri"
                para.alignment = PP_ALIGN.CENTER
                break  # Only process first paragraph

            # Remove companion dashed rectangle if we found one
            if companion is not None:
                try:
                    slide.shapes._spTree.remove(companion._element)
                except Exception:
                    if debug:
                        print(f"  Slide {slide_idx}: could not remove companion rect for '{label}'",
                              file=sys.stderr)

            matched_labels.append(label)
            used_labels.add(norm)

            if debug:
                print(f"  Slide {slide_idx}: replaced '{label}' with {img_rel} "
                      f"({w}x{h} EMU at {x},{y})", file=sys.stderr)

    # Report unused manifest figures
    unused_figures: List[str] = []
    for entry in manifest:
        norm = normalize_label(entry.get("label", ""))
        if norm and norm not in used_labels:
            unused_figures.append(entry.get("label", "?"))

    # Summary
    print(f"\n--- Figure insertion summary ---", file=sys.stderr)
    print(f"  Replaced:              {len(matched_labels)}", file=sys.stderr)
    print(f"  Unmatched placeholders: {len(unmatched_placeholders)}", file=sys.stderr)
    print(f"  Unused manifest figs:  {len(unused_figures)}", file=sys.stderr)

    if matched_labels and debug:
        print(f"  Matched: {', '.join(matched_labels)}", file=sys.stderr)
    if unmatched_placeholders:
        print(f"  Unmatched: {', '.join(unmatched_placeholders)}", file=sys.stderr)
    if unused_figures:
        print(f"  Unused:  {', '.join(unused_figures)}", file=sys.stderr)

    # Save
    prs.save(output_path)
    print(f"\nSaved: {output_path}", file=sys.stderr)

    return True


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Insert extracted figures into a PPTX presentation, "
                    "replacing placeholder shapes with actual images."
    )
    parser.add_argument(
        "pptx_path",
        help="Path to the input PPTX file containing placeholders.",
    )
    parser.add_argument(
        "manifest_json",
        help="Path to the figure manifest JSON file (from extract_figures.py).",
    )
    parser.add_argument(
        "-o", "--output",
        default=None,
        help="Output PPTX path. Defaults to <input>_with_figures.pptx.",
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Enable verbose diagnostic output.",
    )

    args = parser.parse_args()

    # Validate inputs
    if not os.path.isfile(args.pptx_path):
        print(f"Error: PPTX file not found: {args.pptx_path}", file=sys.stderr)
        return 1
    if not os.path.isfile(args.manifest_json):
        print(f"Error: Manifest file not found: {args.manifest_json}", file=sys.stderr)
        return 1

    # Determine output path
    output_path = args.output
    if output_path is None:
        stem = Path(args.pptx_path).stem
        parent = Path(args.pptx_path).parent
        output_path = str(parent / f"{stem}_with_figures.pptx")

    try:
        ok = process_presentation(
            pptx_path=args.pptx_path,
            manifest_path=args.manifest_json,
            output_path=output_path,
            debug=args.debug,
        )
        return 0 if ok else 1
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        if args.debug:
            import traceback
            traceback.print_exc(file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
